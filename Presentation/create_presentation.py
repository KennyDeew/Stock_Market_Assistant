#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Скрипт для создания финальной презентации Stock Market Assistant
на основе шаблона и детального описания слайдов
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
from pathlib import Path
from datetime import datetime
import cairosvg
from io import BytesIO
from PIL import Image

# Цветовая схема
ACCENT_COLOR = RGBColor(56, 189, 248)  # #38bdf8
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)
DARK_BG = RGBColor(30, 30, 30)

def svg_to_png(svg_path, output_path=None, width=1920):
    """Конвертировать SVG в PNG"""
    if not os.path.exists(svg_path):
        return None

    if output_path is None:
        output_path = svg_path.replace('.svg', '.png')

    try:
        # Конвертируем SVG в PNG
        png_data = cairosvg.svg2png(url=str(svg_path), output_width=width)
        with open(output_path, 'wb') as f:
            f.write(png_data)
        return output_path
    except Exception as e:
        print(f"Ошибка конвертации SVG {svg_path}: {e}")
        return None

def add_text_box(slide, left, top, width, height, text, font_size, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    """Добавить текстовое поле на слайд"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = text

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = alignment
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color

    return textbox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE):
    """Добавить маркированный список на слайд"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = ""

    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = item
        p.level = 0
        p.alignment = PP_ALIGN.LEFT

        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

    return textbox

def add_image(slide, image_path, left, top, width=None, height=None):
    """Добавить изображение на слайд"""
    if not os.path.exists(image_path):
        print(f"Warning: Image not found: {image_path}")
        return None

    try:
        if width and height:
            return slide.shapes.add_picture(image_path, Inches(left), Inches(top),
                                           width=Inches(width), height=Inches(height))
        else:
            return slide.shapes.add_picture(image_path, Inches(left), Inches(top))
    except Exception as e:
        print(f"Ошибка вставки изображения {image_path}: {e}")
        return None

def create_presentation():
    """Создать презентацию на основе шаблона"""

    # Пути к файлам
    template_path = Path("Presentation/StockMarketAssistant_Presentation_template.pptx")
    output_path = Path("Presentation/StockMarketAssistant_Presentation_Final.pptx")
    diagrams_path = Path("Presentation/Sequence_Diagrams")

    # Открыть шаблон
    if template_path.exists():
        prs = Presentation(str(template_path))
        print(f"Открыт шаблон: {template_path}")
        # Очистить существующие слайды (кроме первого, если он есть)
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        # Создать новую презентацию, если шаблон не найден
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        print("Создана новая презентация (шаблон не найден)")

    # СЛАЙД 1: ТИТУЛЬНЫЙ СЛАЙД
    slide = prs.slides.add_slide(prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Stock Market Assistant"

    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Платформа для анализа биржевых котировок с многопользовательским доступом"

    # Добавить авторов
    authors_text = "Дубровский Никита Владимирович\nЗаворотный Александр Александрович\nМельников Игорь Евгеньевич\nШадрин Максим Александрович\nПавлов Константин Петрович"
    add_text_box(slide, 1, 4.5, 8, 2, authors_text, 18, False, LIGHT_GRAY, PP_ALIGN.CENTER)

    # Дата
    date_text = datetime.now().strftime("%d.%m.%Y")
    add_text_box(slide, 1, 6.5, 8, 0.5, date_text, 16, False, LIGHT_GRAY, PP_ALIGN.CENTER)

    print("✓ Создан слайд 1: Титульный")

    # СЛАЙД 2: ПРОБЛЕМА И РЕШЕНИЕ
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Проблема и решение"

    # Проблема (слева)
    problem_items = [
        "Инвесторам нужна единая платформа для мониторинга котировок в реальном времени",
        "Существующие решения: сложные, дорогие или ограниченные в функционале",
        "Нет удобного многопользовательского доступа с кастомизацией"
    ]
    add_text_box(slide, 0.5, 1.5, 4, 0.5, "Проблема", 20, True, RGBColor(255, 100, 100))
    add_bullet_list(slide, 0.5, 2, 4, 4, problem_items, 14, WHITE)

    # Решение (справа)
    solution_items = [
        "Веб-платформа с real-time котировками через WebSocket",
        "Простой интерфейс для портфелей и оповещений",
        "Многопользовательский доступ с правами доступа (RBAC)",
        "Аналитика транзакций и рейтинги активов"
    ]
    add_text_box(slide, 5.5, 1.5, 4, 0.5, "Решение", 20, True, RGBColor(100, 255, 100))
    add_bullet_list(slide, 5.5, 2, 4, 4, solution_items, 14, WHITE)

    # Целевая аудитория
    audience_text = "Целевая аудитория: Розничные инвесторы • Профессиональные трейдеры • Финансовые аналитики • Инвестиционные команды"
    add_text_box(slide, 0.5, 6.5, 9, 0.5, audience_text, 14, False, ACCENT_COLOR, PP_ALIGN.CENTER)

    print("✓ Создан слайд 2: Проблема и решение")

    # СЛАЙД 3: ОБЗОР СИСТЕМЫ
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Обзор системы"

    capabilities = [
        "Real-time мониторинг котировок акций, облигаций, криптовалют",
        "Управление инвестиционными портфелями",
        "Трекинг транзакций (покупка/продажа)",
        "Система оповещений (Email) при достижении целевых цен",
        "Аналитика и рейтинги активов",
        "Многопользовательский доступ с разделением прав"
    ]
    add_bullet_list(slide, 0.5, 1.5, 9, 4, capabilities, 16, WHITE)

    asset_types = "Типы активов: Акции (Shares) • Облигации (Bonds) • Криптовалюты (Crypto)"
    add_text_box(slide, 0.5, 6, 9, 0.5, asset_types, 18, True, ACCENT_COLOR, PP_ALIGN.CENTER)

    print("✓ Создан слайд 3: Обзор системы")

    # СЛАЙД 4: ДОМЕННАЯ МОДЕЛЬ
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Доменная модель"

    domain_text = """PortfolioService: Portfolio, PortfolioAsset, PortfolioAssetTransaction
StockCardService: ShareCard, BondCard, CryptoCard, Multiplier, Dividend, Coupon
AnalyticsService: AssetTransaction, AssetRating, Period
AuthService: User, Role, Permission

Ключевые связи:
User → Portfolio (1:N) • Portfolio → PortfolioAsset (1:N) • PortfolioAsset → PortfolioAssetTransaction (1:N)
PortfolioAsset → StockCard (N:1) • PortfolioAssetTransaction → AssetTransaction (1:1 через Kafka)"""

    add_text_box(slide, 0.5, 1.5, 9, 5.5, domain_text, 14, False, WHITE, PP_ALIGN.LEFT)

    print("✓ Создан слайд 4: Доменная модель")

    # СЛАЙД 5: ОСНОВНЫЕ БИЗНЕС-ПРОЦЕССЫ (с диаграммами)
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Основные бизнес-процессы"

    processes = [
        "1. Мониторинг котировок в реальном времени (SignalR, MOEX API)",
        "2. Управление портфелем (CRUD операции, расчет PnL)",
        "3. Система оповещений (Email через NotificationService)",
        "4. Аналитика и рейтинги (Kafka, batch processing, агрегация)",
        "5. Многопользовательский доступ (RBAC, приватные портфели)"
    ]
    add_bullet_list(slide, 0.5, 1.5, 9, 2.5, processes, 14, WHITE)

    # Вставить диаграммы для бизнес-процессов
    # diagram_003.svg - Получение котировок в реальном времени
    # diagram_004.svg - Система оповещений
    # diagram_005.svg - Аналитика

    diagram_files = [
        ("diagram_003.svg", "Котировки в реальном времени", 0.5, 4, 2.8, 2),
        ("diagram_004.svg", "Система оповещений", 3.5, 4, 2.8, 2),
        ("diagram_005.svg", "Аналитика", 6.5, 4, 2.8, 2)
    ]

    for diagram_file, label, left, top, width, height in diagram_files:
        svg_path = diagrams_path / diagram_file
        if svg_path.exists():
            # Конвертируем SVG в PNG
            png_path = svg_path.with_suffix('.png')
            png_path = svg_to_png(svg_path, str(png_path))
            if png_path and os.path.exists(png_path):
                add_image(slide, png_path, left, top, width, height)
                add_text_box(slide, left, top + height + 0.1, width, 0.2, label, 10, False, ACCENT_COLOR, PP_ALIGN.CENTER)
                print(f"  Вставлена диаграмма: {diagram_file}")

    print("✓ Создан слайд 5: Основные бизнес-процессы")

    # СЛАЙД 6: USER STORIES
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "User Stories (Фазы разработки)"

    phases_text = """Phase 1 (MVP — неделя 1-3):
• Real-time quote display
• Portfolio CRUD (Create, Read, Update, Delete)
• Basic transaction tracking
• Authentication и авторизация

Phase 2 (WebSocket & Alerts — неделя 4-6):
• SignalR для real-time quotes
• Email оповещения через NotificationService
• Advanced alert management
• 100 concurrent users

Phase 3 (Analytics — неделя 7-9):
• AnalyticsService: рейтинги активов
• Агрегация данных по периодам
• Топ активов по покупкам/продажам
• 1000 concurrent users

Phase 4 (Production hardening — неделя 10-11):
• Security audit & penetration testing
• Monitoring & observability (OpenTelemetry, OpenSearch)
• Complete documentation
• User acceptance testing"""

    add_text_box(slide, 0.5, 1.5, 9, 5.5, phases_text, 12, False, WHITE, PP_ALIGN.LEFT)

    print("✓ Создан слайд 6: User Stories")

    # СЛАЙД 7: ТЕХНОЛОГИЧЕСКИЙ СТЕК
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Технологический стек"

    tech_stack = """Frontend: React 18 + TypeScript, Redux Toolkit, Tailwind CSS, SignalR client, Vite
Backend: .NET 8+ LTS, ASP.NET Core Web API, C# 12, Entity Framework Core, SignalR, JWT, Autofac, NSwag, Serilog
Database: PostgreSQL 17+, Redis 7+, MongoDB 8+
Message Queue: Apache Kafka, Confluent Kafka .NET client, Kafka UI
Infrastructure: .NET Aspire, Docker + Docker Compose, Kubernetes, GitHub Actions, AWS/Azure/Yandex.Cloud
Monitoring: OpenTelemetry, OpenSearch, OpenSearch Dashboards
External APIs: MOEX (Московская биржа), SendGrid (email)"""

    add_text_box(slide, 0.5, 1.5, 9, 5.5, tech_stack, 12, False, WHITE, PP_ALIGN.LEFT)

    print("✓ Создан слайд 7: Технологический стек")

    # СЛАЙД 8: МИКРОСЕРВИСНАЯ АРХИТЕКТУРА
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Микросервисная архитектура"

    services_text = """6 микросервисов:

1. Gateway Service - API Gateway для маршрутизации запросов, единая точка входа
2. AuthService - Аутентификация (JWT токены), авторизация (RBAC), управление пользователями
3. StockCardService - Управление карточками активов, интеграция с MOEX, real-time котировки через SignalR
4. PortfolioService - Управление портфелями пользователей, регистрация транзакций, публикация событий в Kafka
5. AnalyticsService - Потребление событий транзакций из Kafka, расчет рейтингов активов, агрегация данных
6. NotificationService - Потребление событий из Kafka, отправка Email уведомлений

Архитектурный стиль:
Clean Architecture • Event-Driven Architecture (Kafka) • RESTful API • Database per Service"""

    add_text_box(slide, 0.5, 1.5, 9, 5.5, services_text, 12, False, WHITE, PP_ALIGN.LEFT)

    print("✓ Создан слайд 8: Микросервисная архитектура")

    # СЛАЙД 9: АРХИТЕКТУРА СИСТЕМЫ (LAYERS)
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Архитектура системы (Layers)"

    layers_text = """Clean Architecture Layers:

1. Presentation Layer (WebApi)
   - ASP.NET Core Controllers, DTOs, Middleware, SignalR Hubs

2. Application Layer
   - Use Cases, Application Services, DTOs и маппинг, Валидаторы (FluentValidation)

3. Domain Layer
   - Entities, Value Objects, Domain Services, Domain Events, Enums и константы

4. Infrastructure Layer
   - Entity Framework Core, Repositories, Kafka Consumers/Producers, HTTP Clients
   - External API integrations (MOEX, SendGrid), Caching (Redis)
   - Background Services, Outbox Pattern

Data Flow:
Browser → Gateway → Backend Services → PostgreSQL/Redis/MongoDB
PortfolioService → Kafka → AnalyticsService/NotificationService
SignalR → Real-time quote broadcasting"""

    add_text_box(slide, 0.5, 1.5, 9, 5.5, layers_text, 11, False, WHITE, PP_ALIGN.LEFT)

    print("✓ Создан слайд 9: Архитектура системы")

    # СЛАЙД 10: DATABASE SCHEMA
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Database Schema (ER-модель)"

    db_text = """PortfolioService (PostgreSQL):
portfolio, portfolio_asset, portfolio_asset_transaction

StockCardService (PostgreSQL + MongoDB):
ShareCards, BondCards, CryptoCards, Multipliers, Dividends, Coupons
MongoDB: FinancialReports (коллекция)

AnalyticsService (PostgreSQL):
asset_transactions, asset_ratings

AuthService (PostgreSQL):
users, roles, role_permissions, permissions, refresh_sessions

Key Optimizations:
Индексы на внешние ключи (portfolio_id, user_id, stock_card_id)
Индексы на даты (transaction_date, period_start, period_end)
Индексы для поиска (ticker, name)
Кэширование часто запрашиваемых данных в Redis"""

    add_text_box(slide, 0.5, 1.5, 9, 5.5, db_text, 11, False, WHITE, PP_ALIGN.LEFT)

    print("✓ Создан слайд 10: Database Schema")

    # СЛАЙД 11: NON-FUNCTIONAL REQUIREMENTS
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Non-Functional Requirements (NFR)"

    nfr_text = """Производительность:
API response time: p95 < 200ms, p99 < 500ms • WebSocket latency: < 100ms
Database query time: < 50ms (p95) • Поддержка 1000+ concurrent users • Throughput: 1000+ requests/second

Масштабируемость:
Горизонтальное масштабирование микросервисов • Автоматическое масштабирование в Kubernetes (HPA)
Кэширование через Redis • Асинхронная обработка через Kafka • Database per Service

Надежность:
Availability: 99.9% (SLA) • Retry механизмы (Polly) • Circuit Breaker • Graceful degradation • Health checks

Безопасность:
TLS 1.3 • JWT токены (1 час) • Refresh token rotation • Rate limiting (100 req/min)
Валидация входных данных (FluentValidation) • SQL injection protection (EF Core) • CORS restricted

Поддерживаемость:
Clean Architecture • Unit test coverage: 80%+ • Integration tests (TestContainers)
Структурированное логирование (Serilog + OpenSearch) • Мониторинг (OpenTelemetry) • API documentation (Swagger)"""

    add_text_box(slide, 0.5, 1.5, 9, 5.5, nfr_text, 10, False, WHITE, PP_ALIGN.LEFT)

    print("✓ Создан слайд 11: Non-Functional Requirements")

    # СЛАЙД 12: SECURITY ARCHITECTURE
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Security Architecture"

    security_text = """Authentication:
JWT tokens с коротким временем жизни (15 минут для access token)
Refresh token rotation (автоматическая смена, 30 дней)
Password policy: min 8 chars, complexity requirements
ASP.NET Identity для управления пользователями

Authorization:
RBAC: ADMIN, USER roles • Row-level security для user data (фильтрация по user_id)
Policy-based authorization в ASP.NET Core • Приватные портфели (IsPrivate флаг)

Data Protection:
TLS 1.3 for transit (HTTPS) • AES-256 for data at rest
Password hashing (bcrypt через ASP.NET Identity) • Secrets management (environment variables, Azure Key Vault)

API Security:
Rate limiting per user (100 req/min) • CORS restricted to registered domains
API key validation для внешних интеграций • Request signing для sensitive operations
Input validation (FluentValidation) • SQL injection protection (EF Core)"""

    add_text_box(slide, 0.5, 1.5, 9, 5.5, security_text, 11, False, WHITE, PP_ALIGN.LEFT)

    print("✓ Создан слайд 12: Security Architecture")

    # СЛАЙД 13: TESTING STRATEGY
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Testing Strategy"

    testing_text = """Unit Tests (80% coverage):
xUnit / NUnit (test framework) • Moq / NSubstitute (mocking) • FluentAssertions (assertions)
Domain services, application services, utilities • Component logic tests

Integration Tests:
API endpoints с database (TestContainers для PostgreSQL)
External API mocking (MOEX, SendGrid) • Kafka integration tests (Testcontainers)
HTTP client testing (WebApplicationFactory)

Performance Tests:
Load testing (k6, NBomber) • Stress testing для определения пределов
Database performance tests

E2E Tests:
Playwright / Cypress для frontend • API E2E tests через Gateway"""

    add_text_box(slide, 0.5, 1.5, 9, 5.5, testing_text, 11, False, WHITE, PP_ALIGN.LEFT)

    print("✓ Создан слайд 13: Testing Strategy")

    # СЛАЙД 14: DEPLOYMENT & DEVOPS
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Deployment & DevOps"

    devops_text = """CI/CD Pipeline (GitHub Actions):
1. Code push to Git 2. Linting & formatting checks 3. Unit tests 4. Integration tests
5. Build Docker image 6. Push to container registry 7. Deploy to staging
8. Run E2E tests 9. Manual approval 10. Deploy to production (blue-green deployment)

Infrastructure:
Docker containers для всех сервисов • .NET Aspire для локальной разработки
Kubernetes для orchestration (production) • Database: Managed PostgreSQL (AWS RDS / Azure Database)
Cache: Redis cluster (managed service) • Load balancer: ALB (AWS) / Azure LB
CDN: Cloudflare / AWS CloudFront • Monitoring: OpenTelemetry • Logging: OpenSearch
Dashboards: OpenSearch Dashboards • Kafka UI (порт 9100) • Mongo Express (порт 5005) • PgWeb (порты 5000, 5001)

Deployment Strategy:
Blue-Green deployment для zero-downtime • Canary releases для постепенного rollout
Database migrations через EF Core Migrations (автоматические при старте)
Health checks для всех сервисов (/health, /alive)"""

    add_text_box(slide, 0.5, 1.5, 9, 5.5, devops_text, 10, False, WHITE, PP_ALIGN.LEFT)

    print("✓ Создан слайд 14: Deployment & DevOps")

    # СЛАЙД 15: MONITORING & OBSERVABILITY
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Monitoring & Observability"

    monitoring_text = """Metrics & Tracing (OpenTelemetry):
API response time (p50, p95, p99) • Database query time • Cache hit ratio (Redis)
WebSocket connections (SignalR) • Kafka message processing rate
Error rates по сервисам • Distributed tracing между сервисами

Dashboards (OpenSearch Dashboards):
System health (CPU, memory, disk) • API performance (latency, throughput)
Database metrics (connections, slow queries) • Business metrics (active users, transactions, portfolios)
Error rates и logs • Kafka topics monitoring

Logging (OpenSearch):
Structured JSON logs • Log levels: ERROR, WARN, INFO, DEBUG
Centralized search и analysis • Correlation IDs для трейсинга запросов
Интеграция с OpenTelemetry для контекста

Observability:
OpenTelemetry для метрик, трейсинга и логов
OpenSearch для хранения и анализа логов • OpenSearch Dashboards для визуализации"""

    add_text_box(slide, 0.5, 1.5, 9, 5.5, monitoring_text, 10, False, WHITE, PP_ALIGN.LEFT)

    print("✓ Создан слайд 15: Monitoring & Observability")

    # СЛАЙД 16: EVENT-DRIVEN COMMUNICATION (с диаграммами)
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Event-Driven Communication"

    kafka_text = """Kafka Topics:

1. portfolio.transactions
   - Producer: PortfolioService
   - Consumers: AnalyticsService, NotificationService
   - Consumer Group: analytics-service-transactions
   - Batch processing: 100 сообщений за раз
   - Manual offset commit после успешной обработки

2. financial.report.created
   - Producer: StockCardService
   - Payload: FinancialReportCreatedMessage
   - Хранение отчетов в MongoDB

Схема взаимодействия:
PortfolioService → Outbox → Kafka → AnalyticsService/NotificationService"""

    add_text_box(slide, 0.5, 1.5, 4.5, 4, kafka_text, 11, False, WHITE, PP_ALIGN.LEFT)

    # Вставить диаграмму Kafka (diagram_002.svg - Создание транзакции с публикацией в Kafka)
    svg_path = diagrams_path / "diagram_002.svg"
    if svg_path.exists():
        png_path = svg_path.with_suffix('.png')
        png_path = svg_to_png(svg_path, str(png_path))
        if png_path and os.path.exists(png_path):
            add_image(slide, png_path, 5.5, 1.5, 4.5, 4)
            print(f"  Вставлена диаграмма: diagram_002.svg")

    print("✓ Создан слайд 16: Event-Driven Communication")

    # СЛАЙД 17: API ENDPOINTS
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "API Endpoints"

    api_text = """PortfolioService:
GET/POST /api/v1/portfolios • GET/PUT/DELETE /api/v1/portfolios/{id}
GET/POST /api/v1/portfolio-assets • GET/POST /api/v1/portfolio-assets/{id}/transactions

StockCardService:
GET /api/stockcard/shares • GET /api/stockcard/bonds • GET /api/stockcard/crypto
GET /api/stockcard/{ticker} • GET /api/stockcard/{ticker}/price
SignalR Hub: /priceHub (подписка на real-time котировки)

AnalyticsService:
GET /api/analytics/transactions • GET /api/analytics/assets/top-bought
GET /api/analytics/assets/top-sold • GET /api/analytics/portfolios/{id}/history
POST /api/analytics/portfolios/compare

AuthService:
POST /api/auth/register • POST /api/auth/login • POST /api/auth/refresh • POST /api/auth/logout

NotificationService:
POST /api/notifications/send • Потребление событий из Kafka для автоматических уведомлений"""

    add_text_box(slide, 0.5, 1.5, 9, 5.5, api_text, 10, False, WHITE, PP_ALIGN.LEFT)

    print("✓ Создан слайд 17: API Endpoints")

    # СЛАЙД 18: ВЫВОДЫ И ПЕРСПЕКТИВЫ
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[6])

    if slide.shapes.title:
        slide.shapes.title.text = "Выводы и перспективы"

    results_text = """Технические результаты:
✓ Масштабируемая микросервисная архитектура с Clean Architecture
✓ Real-time infrastructure с SignalR WebSocket
✓ Event-Driven Architecture с Kafka
✓ Оптимизированные databases с индексами
✓ Multi-layer security (JWT, encryption)
✓ Comprehensive monitoring и observability

Бизнес-результаты:
✓ Единая платформа для мониторинга котировок
✓ Удобное управление портфелями
✓ Автоматические оповещения
✓ Аналитика и рейтинги активов
✓ Многопользовательский доступ с RBAC

Перспективы:
1. Масштабирование в Kubernetes для production
2. Использование ИИ для аналитики и прогнозирования
3. Расширение типов активов (фьючерсы, опционы)
4. Мобильное приложение (React Native)
5. Расширенная аналитика с ML моделями"""

    add_text_box(slide, 0.5, 1.5, 9, 5.5, results_text, 11, False, WHITE, PP_ALIGN.LEFT)

    print("✓ Создан слайд 18: Выводы и перспективы")

    # Сохранить презентацию
    prs.save(str(output_path))
    print(f"\n{'='*60}")
    print(f"Презентация сохранена: {output_path}")
    print(f"Всего слайдов: {len(prs.slides)}")
    print(f"{'='*60}")

if __name__ == "__main__":
    create_presentation()
